from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import requests
from requests.exceptions import HTTPError
from bs4 import BeautifulSoup
import spacy
import os
import datetime
from dotenv import load_dotenv
load_dotenv()

PIXABAY_API_KEY = os.getenv('PIXABAY_API_KEY')
IMAGE_TYPE_TO_SEARCH = os.getenv('IMAGE_TYPE_TO_SEARCH')
SOURCE_TXT_FILE = os.getenv('SOURCE_TXT_FILE')
STORY_URL_TO_DOWNLOAD = os.getenv('STORY_URL_TO_DOWNLOAD')
NUMBER_OF_SLIDES_PER_PRESENTATION = int(os.getenv(
    'NUMBER_OF_SLIDES_PER_PRESENTATION'))
NUMBER_OF_WORDS = int(os.getenv('NUMBER_OF_WORDS'))


nlp = spacy.load("es_core_news_sm")
img_left = img_top = Inches(0)


def get_text_from_archive_org(local_file):
    file_url = STORY_URL_TO_DOWNLOAD
    html_content = requests.get(file_url).text
    # Parse the html content
    soup = BeautifulSoup(html_content, "lxml")
    pre_tag = soup.find("pre")
    open(local_file, 'w').write(pre_tag.text)


def get_text_from_gutenberg_org(local_file):
    file_url = STORY_URL_TO_DOWNLOAD
    html_content = requests.get(file_url).text
    open(local_file, 'w').write(html_content)


def process_with_spacy(local_file):
    with open(local_file, "r") as file:
        text_file = file.read()
    doc = nlp(text_file)

    nouns = filter_nouns(doc)
    print(f'====== len(nouns): {len(nouns)} =====')
    write_to_presentation(
        nouns, f'doman_method_nouns_{datetime.datetime.now().strftime("%Y-%m-%d")}.pptx', NUMBER_OF_SLIDES_PER_PRESENTATION, NUMBER_OF_WORDS)

    verbs = filter_verbs(doc)
    print(f'===== len(verbs): {len(verbs)} =====')
    write_to_presentation(
        verbs, f'doman_method_verbs_{datetime.datetime.now().strftime("%Y-%m-%d")}.pptx', NUMBER_OF_SLIDES_PER_PRESENTATION, NUMBER_OF_WORDS)

    # Find named entities, phrases and concepts
    entities = filter_entities(doc)
    print(f'===== len(entities): {len(entities)} =====')
    write_to_presentation(
        entities, f'doman_method_entities_{datetime.datetime.now().strftime("%Y-%m-%d")}.pptx', NUMBER_OF_SLIDES_PER_PRESENTATION, NUMBER_OF_WORDS)


def write_to_presentation(data, doman_file, number_of_slides, number_of_words):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    cont = 0
    for word in data:
        if len(word.strip()) > 3 and len(word.strip().split()) == number_of_words:
            slide = prs.slides.add_slide(title_slide_layout)

            title = slide.shapes.title

            tf = title.text_frame

            # paragrap in the text frame
            p = tf.add_paragraph()
            p.text = word.lower()
            p.font.bold = True
            p.font.name = 'Arial'
            p.font.size = Pt(70)
            p.font.color.rgb = RGBColor(255, 0, 0)

            image_url = get_url_image(word.lower())
            print(f'image_url: {image_url}')
            if image_url:
                image_file_name = '/tmp/' + image_url.split('/')[-1]
                save_image(image_url, image_file_name)
                print(f'image_local_path: {image_file_name}')

                slide = prs.slides.add_slide(title_slide_layout)
                pic = slide.shapes.add_picture(
                    image_file_name, img_left, img_top)
            cont += 1
        else:
            print(f'No recorded: {word}')
        if cont == number_of_slides:
            break
    prs.save(doman_file)


def get_url_image(word_searched):
    api_url = "https://pixabay.com/api/"
    search_params = {
        'key': PIXABAY_API_KEY,
        'q': word_searched,
        'image_type': IMAGE_TYPE_TO_SEARCH,
        'safesearch': 'true',
        'order': 'ec',
        'per_page': '3'
    }
    try:
        response = requests.get(api_url, params=search_params)

        # If the response was successful, no Exception will be raised
        response.raise_for_status()
    except HTTPError as http_err:
        print(f'HTTP error occurred: {http_err}')  # Python 3.6
    except Exception as err:
        print(f'Other error occurred: {err}')  # Python 3.6
    else:
        print('Success!')
        print(f'Word: {word_searched}')
        print(f'response: {response}')
        dict_response = response.json()
        # print(f'dict_response: {dict_response}')
        if dict_response['total'] != 0:
            return (dict_response["hits"][0]["webformatURL"])
    return ""


def save_image(image_url, image_file_name):
    with open(image_file_name, "wb") as file:
        response = requests.get(image_url)
        file.write(response.content)


def filter_entities(doc):
    print('=========== Named entities, phrases and concepts ======')
    set_entities = set()
    for entity in doc.ents:
        cleaned_phrase = entity.text.strip().replace('\n', '')
        if entity.label_ in ['LOC', 'ORG'] and cleaned_phrase and len(cleaned_phrase) > 1 and '«' not in cleaned_phrase:
            set_entities.add(cleaned_phrase)
    return set_entities


def filter_verbs(doc):
    print('=========== Verbs ======')
    set_verbs = set()
    for token in doc:
        if token.pos_ == "VERB":
            set_verbs.add(token.lemma_)
    return set_verbs


def filter_nouns(doc):
    print('=========== Nouns ======')
    set_nouns = set()
    for chunk in doc.noun_chunks:
        set_nouns.add(chunk.text)
    return set_nouns


def read_txt_from_archive_org(local_file):
    get_text_from_archive_org(local_file)
    print(f'File Size is {os.stat(local_file).st_size / (1024 * 1024)} MB')

    process_with_spacy(local_file)


def read_txt_from_gutenberg_org(local_file):
    get_text_from_gutenberg_org(local_file)
    print(f'File Size is: {os.stat(local_file).st_size / (1024 * 1024)} MB')

    process_with_spacy(local_file)


if __name__ == "__main__":
    local_file = "/tmp/" + STORY_URL_TO_DOWNLOAD.split('/')[-1]
    if os.path.exists(local_file):
        print(f'Local File: {os.stat(local_file).st_size / (1024 * 1024)} MB')
        process_with_spacy(local_file)
    elif SOURCE_TXT_FILE == 'archive.org':
        print("Leyendo de archive.org")
        read_txt_from_archive_org(local_file)
    elif SOURCE_TXT_FILE == 'gutenberg.org':
        print("Leyendo de gutenberg.org")
        read_txt_from_gutenberg_org(local_file)
    else:
        print("Algo salió mal")
